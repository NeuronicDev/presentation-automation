import sys
import json
import logging
import traceback
from io import StringIO
import os

try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT, MSO_TEXT_UNDERLINE_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE, MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE, PP_MEDIA_TYPE
    from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE, MSO_COLOR_TYPE, MSO_PATTERN_TYPE, MSO_THEME_COLOR_INDEX
    from pptx.chart.data import CategoryChartData, ChartData, XyChartData, BubbleChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE, XL_DATA_LABEL_POSITION
    from pptx.enum.action import PP_ACTION_TYPE
    from pptx.table import Table, _Cell 

except ImportError as import_err:
    print(json.dumps({"status": "failed", "errors": [{"error": f"Import error in container: {import_err}. Ensure Dockerfile installs all dependencies."}]}), file=sys.stderr)
    sys.exit(1)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - CONTAINER - %(levelname)s - %(message)s', stream=sys.stderr) 

PPTX_CONTAINER_PATH = "/app/presentation.pptx"

def execute_code_snippet(code_snippet: str, prs: Presentation, slide_index: int | None, params: dict) -> tuple[bool, str]:
    log_stream = StringIO() 
    original_stdout = sys.stdout
    original_stderr = sys.stderr
    
    sys.stdout = log_stream
    sys.stderr = log_stream

    success = False
    message = ""
    slide = None 

    try:
        exec_globals = {
            
            # Core Modules/Classes
            'pptx': pptx,
            'Presentation': Presentation,
            'Table': Table,
            '_Cell': _Cell, 
            'ChartData': ChartData,
            'CategoryChartData': CategoryChartData,
            'XyChartData': XyChartData,
            'BubbleChartData': BubbleChartData, 
            'PP_ACTION_TYPE': PP_ACTION_TYPE,
            'PP_PLACEHOLDER_TYPE': PP_PLACEHOLDER_TYPE,
            'PP_MEDIA_TYPE': PP_MEDIA_TYPE,
            
            # Utilities
            'Inches': Inches,
            'Pt': Pt,
            'Emu': Emu,
            'Cm': Cm,
            'RGBColor': RGBColor,

            # Text Enums
            'MSO_AUTO_SIZE': MSO_AUTO_SIZE,
            'MSO_VERTICAL_ANCHOR': MSO_VERTICAL_ANCHOR,
            'PP_PARAGRAPH_ALIGNMENT': PP_PARAGRAPH_ALIGNMENT,
            'MSO_TEXT_UNDERLINE_TYPE': MSO_TEXT_UNDERLINE_TYPE,

            # Shape Enums
            'MSO_SHAPE_TYPE': MSO_SHAPE_TYPE,
            'MSO_CONNECTOR_TYPE': MSO_CONNECTOR_TYPE,
            'MSO_AUTO_SHAPE_TYPE': MSO_AUTO_SHAPE_TYPE,
            'PP_PLACEHOLDER_TYPE': PP_PLACEHOLDER_TYPE,

            # Fill/Line/Effect Enums
            'MSO_FILL_TYPE': MSO_FILL_TYPE,
            'MSO_THEME_COLOR_INDEX': MSO_THEME_COLOR_INDEX,
            'MSO_COLOR_TYPE': MSO_COLOR_TYPE,
            'MSO_PATTERN_TYPE': MSO_PATTERN_TYPE,
            'MSO_LINE_DASH_STYLE': MSO_LINE_DASH_STYLE,

            # Chart Enums
            'XL_CHART_TYPE': XL_CHART_TYPE,
            'XL_LEGEND_POSITION': XL_LEGEND_POSITION,
            'XL_MARKER_STYLE': XL_MARKER_STYLE,
            'XL_DATA_LABEL_POSITION': XL_DATA_LABEL_POSITION,
            'XL_TICK_MARK': XL_TICK_MARK,
            'XL_TICK_LABEL_POSITION': XL_TICK_LABEL_POSITION,

            # Standard Builtins
            '__builtins__': __builtins__
        }


        exec_locals = {
            'prs': prs, 
            'slide': None,
            'params': params if params else {}
        }

        if slide_index is not None:
            if 0 <= slide_index < len(prs.slides):
                exec_locals['slide'] = prs.slides[slide_index]
                logging.info(f"Targeting slide index: {slide_index}")
            else:
                raise IndexError(f"Invalid slide_index {slide_index} received (0-based). Presentation has {len(prs.slides)} slides.")
        else:
             logging.info("No specific slide index provided.")

        compiled_code = compile(code_snippet, '<string>', 'exec')
        exec(compiled_code, exec_globals, exec_locals)

        success = True
        message = f"Successfully executed snippet."
        logging.info(message)

    except Exception as e:
        success = False
        error_trace = traceback.format_exc()
        message = f"Error executing snippet: {e}\nTraceback:\n{error_trace}"
        logging.error(message)
        
    finally:
        sys.stdout = original_stdout
        sys.stderr = original_stderr
        
        captured_output = log_stream.getvalue()
        if captured_output:
             message += f"\n--- Snippet Output ---\n{captured_output}\n----------------------"
        log_stream.close()

    return success, message


def main_container_logic():
    logging.info("Executor script started inside container.")
    results = {"status": "unknown", "errors": [], "processed_count": 0, "success_count": 0}
    any_snippet_failed = False

    try:
        logging.info("Reading input data...")
        input_data = os.environ.get('TASKS_INPUT', '')
        # input_data = sys.stdin.read()
        logging.info(f"Received input data from env: {input_data}...")
        if not input_data:
            raise ValueError("No input data received from env .")

        try:
            tasks = json.loads(input_data)
            logging.info(f"Received {len(tasks)} tasks to execute.")
            if not isinstance(tasks, list):
                 raise ValueError("Input data is not a JSON list.")
        except json.JSONDecodeError as json_err:
            raise ValueError(f"Failed to parse JSON input from stdin: {json_err}. Input received: '{input_data[:500]}...'")

        if not os.path.exists(PPTX_CONTAINER_PATH):
             raise FileNotFoundError(f"Presentation file not found at mounted path: {PPTX_CONTAINER_PATH}")

        logging.info(f"Loading presentation from {PPTX_CONTAINER_PATH}")
        prs = Presentation(PPTX_CONTAINER_PATH)

        for i, task in enumerate(tasks):
            task_index_for_log = i + 1 
            logging.info(f"--- Processing Task {task_index_for_log}/{len(tasks)} ---")

            code = task.get("generated_code")
            slide_number = task.get("slide_number")
            params = task.get("params", {})
            desc = task.get("description", "No description provided")
            original_instr = task.get("original_instruction", "N/A")

            if not code:
                logging.warning(f"Task {task_index_for_log} has no code snippet. Skipping.")
                results["errors"].append({
                    "task_index": i,
                    "error": "No code snippet provided.",
                    "description": desc,
                    "original_instruction": original_instr
                })
                any_snippet_failed = True 
                continue

            slide_index = (slide_number - 1) if slide_number is not None else None
            logging.info(f"Executing code for slide number {slide_number} (index {slide_index}): {desc[:100]}...")

            success, message = execute_code_snippet(code, prs, slide_index, params)
            results["processed_count"] += 1
            if success:
                results["success_count"] += 1 
            else:
                any_snippet_failed = True
                results["errors"].append({
                    "task_index": i,
                    "error": message,
                    "description": desc,
                    "original_instruction": original_instr
                })

            logging.info(f"--- Finished Task {task_index_for_log} ---")

        if results["processed_count"] > 0:
             logging.info(f"Saving modified presentation back to {PPTX_CONTAINER_PATH}")
             prs.save(PPTX_CONTAINER_PATH)
             logging.info("Presentation saved.")
        else:
             logging.info("No tasks processed, presentation not saved.")

        if any_snippet_failed:
            results["status"] = "partial_success" 
            logging.error(f"Execution finished with errors. {len(results['errors'])} snippet(s) failed.")
        elif results["processed_count"] > 0:
             results["status"] = "success"
             logging.info("Execution finished successfully.")
        else:
             results["status"] = "no_tasks_processed"
             logging.info("Execution finished, but no tasks had code to execute.")

    except FileNotFoundError as fnf_err:
        results["status"] = "failed"
        err_msg = f"File system error in container: {fnf_err}"
        results["errors"].append({"task_index": -1, "error": err_msg})
        logging.error(err_msg)
        any_snippet_failed = True 
        
    except ValueError as val_err: 
         results["status"] = "failed"
         err_msg = f"Input data error in container: {val_err}"
         results["errors"].append({"task_index": -1, "error": err_msg})
         logging.error(err_msg)
         any_snippet_failed = True
         
    except Exception as e:
        results["status"] = "failed"
        err_msg = f"An unexpected error occurred in the container's main logic: {e}\n{traceback.format_exc()}"
        results["errors"].append({"task_index": -1, "error": err_msg})
        logging.error(err_msg)
        any_snippet_failed = True
        
    finally:
        print(json.dumps(results, indent=None).strip())
        sys.exit(1 if results["status"] == "failed" else 0)

if __name__ == "__main__":
    main_container_logic()