# # pptx_handler.py 
# from fastapi import APIRouter, HTTPException, Body, status
# from pydantic import BaseModel, Field
# import base64
# import os
# import aiofiles
# import logging

# router = APIRouter()

# # --- Configuration ---
# SAVE_DIR = "./uploaded_pptx"
# os.makedirs(SAVE_DIR, exist_ok=True)

# # --- Pydantic Model for Request Body ---
# class PPTXPayload(BaseModel):
#     base64: str = Field(..., description="Base64 encoded content of the PPTX file")
#     filename: str = "presentation.pptx"

# # --- Logging Setup ---
# logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
# logger = logging.getLogger(__name__)

# # --- PPTX Upload and Conversion Endpoint ---
# @router.post("", status_code=status.HTTP_200_OK, response_model=dict)
# async def upload_pptx(payload: PPTXPayload = Body(...)):
#     """
#     Receives a Base64 encoded PPTX file, saves it, converts it to slide
#     """
#     try:
#         safe_filename = os.path.basename(payload.filename)
#         pptx_path = os.path.join(SAVE_DIR, safe_filename)
#         pptx_bytes = base64.b64decode(payload.base64)

#         async with aiofiles.open(pptx_path, "wb") as f:
#             await f.write(pptx_bytes)

#         logger.info(f"PPTX file saved at {pptx_path}")
#         return { "status": "success", "message": f"File saved as {safe_filename}" }

#     except base64.binascii.Error as e:
#         logger.error("Base64 decoding error", exc_info=True)
#         raise HTTPException(status_code=400, detail="Invalid Base64 data")
#     except Exception as e:
#         logger.error("Unexpected error while saving PPTX", exc_info=True)
#         raise HTTPException(status_code=500, detail=f"Server error: {e}")


from fastapi import APIRouter, HTTPException, Body, status
from pydantic import BaseModel, Field
import base64
import os
import aiofiles
import logging
from pptx import Presentation

from utils.utils import convert_pptx_to_pdf, generate_slide_context

router = APIRouter()

# --- Configuration ---
SAVE_DIR = "./uploaded_pptx"
os.makedirs(SAVE_DIR, exist_ok=True)
SLIDE_IMAGE_DIR = "./uploaded_pptx/slide_images"

# --- Pydantic Model ---
class PPTXPayload(BaseModel):
    base64: str = Field(..., description="Base64 encoded content of the PPTX file")
    filename: str = "presentation.pptx"

# --- Logging ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- Upload + Process Route ---
@router.post("", status_code=status.HTTP_200_OK, response_model=dict)
async def upload_pptx(payload: PPTXPayload = Body(...)):
    """
    Receives a Base64-encoded PPTX file, saves it, then processes it:
    - Converts to PDF
    - Generates slide image + XML for each slide
    - Saves image bytes in a .txt file for each slide
    """
    try:
        safe_filename = os.path.basename(payload.filename)
        pptx_path = os.path.join(SAVE_DIR, safe_filename)
        pptx_bytes = base64.b64decode(payload.base64)

        async with aiofiles.open(pptx_path, "wb") as f:
            await f.write(pptx_bytes)
        logger.info(f"PPTX file saved at {pptx_path}")

        # Prepare output folder
        pptx_name = os.path.splitext(safe_filename)[0]
        slide_dir = os.path.join(SLIDE_IMAGE_DIR, pptx_name)
        pdf_output_dir = os.path.join(slide_dir, "converted_pdfs")
        os.makedirs(pdf_output_dir, exist_ok=True)

        # Step 1: Convert to PDF
        pdf_path = convert_pptx_to_pdf(pptx_path, output_dir=pdf_output_dir)

        # Step 2: Load presentation
        prs = Presentation(pptx_path)

        # Step 3: Generate context for each slide
        for idx, _ in enumerate(prs.slides):
            slide_number = idx + 1
            logger.info(f"Processing slide {slide_number}...")
            generate_slide_context(prs, slide_number, pdf_path, slide_dir)

        return {
            "status": "success",
            "message": f"File saved and processed: {safe_filename}",
            "slides_processed": len(prs.slides)
        }

    except base64.binascii.Error:
        logger.error("Base64 decoding error", exc_info=True)
        raise HTTPException(status_code=400, detail="Invalid Base64 data")
    except FileNotFoundError as e:
        logger.error("File conversion error", exc_info=True)
        raise HTTPException(status_code=500, detail=f"PDF not generated: {e}")
    except Exception as e:
        logger.error("Unexpected error while saving or processing PPTX", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Server error: {e}")

