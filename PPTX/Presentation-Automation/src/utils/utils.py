import os, logging, zipfile, subprocess, base64, pptx
from io import BytesIO
from typing import Tuple, Dict
from pptx import Presentation
from lxml import etree
from pdf2image import convert_from_path
import xml.dom.minidom
from config.config import LIBREOFFICE_PATH

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

def convert_pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    logger = logging.getLogger(__name__)
    os.makedirs(output_dir, exist_ok=True)

    # Convert to absolute paths to avoid LibreOffice path resolution issues
    abs_pptx_path = os.path.abspath(pptx_path)
    abs_output_dir = os.path.abspath(output_dir)

    command = [
        "soffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", abs_output_dir,
        abs_pptx_path
    ]

    try:
        logger.info(f"Running command: {' '.join(command)}")
        result = subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        logger.info(f"LibreOffice stdout:\n{result.stdout.decode()}")
        logger.info(f"LibreOffice stderr:\n{result.stderr.decode()}")

        # Find the resulting PDF file
        pdf_files = [f for f in os.listdir(abs_output_dir) if f.lower().endswith(".pdf")]
        if not pdf_files:
            raise FileNotFoundError(f"No PDF file found in {output_dir} after conversion")

        pdf_path = os.path.join(abs_output_dir, pdf_files[0])
        logger.info(f"PDF conversion successful: {pdf_path}")
        return pdf_path

    except subprocess.CalledProcessError as e:
        logger.error("LibreOffice command failed", exc_info=True)
        raise RuntimeError("LibreOffice failed to convert PPTX to PDF")

def extract_slide_xml_from_ppt(pptx_path: str, slide_number: int) -> str:
    try:
        slide_filename = f"ppt/slides/slide{slide_number}.xml"
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            if slide_filename not in zip_ref.namelist():
                raise FileNotFoundError(f"Slide XML not found: {slide_filename}")

            xml_content = zip_ref.read(slide_filename)
            dom = xml.dom.minidom.parseString(xml_content)
            pretty_xml = dom.toprettyxml()

            logger.info(f"Extracted XML for slide {slide_number}")
            return pretty_xml
    except Exception as e:
        logger.exception(f"Failed to extract XML for slide {slide_number}")
        raise

def extract_slide_xml(prs: Presentation, slide_index: int) -> str:
    try:
        slide = prs.slides[slide_index]
        slide_element = slide._element
        xml = etree.tostring(slide_element, encoding='unicode', pretty_print=True, method='xml')
        logger.info(f"Extracted XML using lxml for slide index {slide_index}")
        return xml
    except Exception as e:
        logger.exception(f"Error extracting XML using lxml for slide {slide_index}")
        raise

def generate_slide_image(pdf_path: str, slide_index: int) -> Tuple[str, bytes]:
    try:
        images = convert_from_path(pdf_path, first_page=slide_index + 1, last_page=slide_index + 1)
        image = images[0]
        buffered = BytesIO()
        image.save(buffered, format="PNG")
        img_bytes = buffered.getvalue()
        base64_image = base64.b64encode(img_bytes).decode('utf-8')
        return base64_image, img_bytes
    except Exception as e:
        logger.exception(f"Error generating image for slide index {slide_index}")
        raise

def generate_slide_context(prs: Presentation, slide_number: int, pdf_path: str, output_dir: str) -> Dict:
    try:
        slide_index = slide_number 
        os.makedirs(output_dir, exist_ok=True)

        # Generate image
        base64_image, img_bytes = generate_slide_image(pdf_path, slide_index)
        image_path = os.path.join(output_dir, f"slide{slide_number}.png")
        with open(image_path, "wb") as f:
            f.write(img_bytes)
        logger.info(f"Saved image for slide {slide_number} to {image_path}")

        # Save the base64 image data as a .txt file
        txt_file_path = os.path.join(output_dir, f"slide{slide_number}_image.txt")
        with open(txt_file_path, "w", encoding="utf-8") as f:
            f.write(base64_image)
        logger.info(f"Saved base64 image data for slide {slide_number} to {txt_file_path}")

        # Generate XML
        xml_string = extract_slide_xml(prs, slide_index)
        xml_path = os.path.join(output_dir, f"slide{slide_number}.xml")
        with open(xml_path, "w", encoding="utf-8") as f:
            f.write(xml_string)
        logger.info(f"Saved XML for slide {slide_number} to {xml_path}")

        return {
            "slide_xml_structure": xml_string,
            "slide_image_base64": base64_image,
            "slide_image_bytes": img_bytes
        }
    except Exception as e:
        logger.error(f"Failed to generate context for slide {slide_number}")
        raise


def update_slide_context_and_save_to_disk(
    pptx_path: str,
    base_output_dir: str
) -> Tuple[Dict[int, Dict], str]:
    try:
        pptx_filename = os.path.splitext(os.path.basename(pptx_path))[0]
        slide_dir = os.path.join(base_output_dir, pptx_filename)
        os.makedirs(slide_dir, exist_ok=True)

        pdf_path = convert_pptx_to_pdf(pptx_path, output_dir=slide_dir)
        prs = pptx.Presentation(pptx_path)

        slide_context_cache = {}

        for idx, _ in enumerate(prs.slides):
            slide_number = idx + 1
            context = generate_slide_context(prs, slide_number, pdf_path, slide_dir)
            slide_context_cache[slide_number] = context
            logger.info(f"Context saved for slide {slide_number}")

        return slide_context_cache, pdf_path
    except Exception as e:
        logger.exception("Failed to update slide context and save to disk.")
        raise