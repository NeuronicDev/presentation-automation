import os
import base64
import tempfile
from pptx import Presentation
from pptx.util import Inches
from openai import OpenAI
from PIL import Image
from io import BytesIO
import logging

client = OpenAI(
    api_key="AIzaSyAv0sTw83EOKcJtoSyT9ug4cnzwGagkMJY",
    base_url="https://generativelanguage.googleapis.com/v1beta/openai/"
)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class PowerPointProcessor:
    def __init__(self, pptx_file):
        self.pptx_file = pptx_file
        self.prs = Presentation(pptx_file)
        self.inserted_slides_count = 0
        self.temp_dir = tempfile.mkdtemp()

    def extract_images_from_pptx(self):
        image_paths = []
        try:
            for slide_num, slide in enumerate(self.prs.slides, start=1):
                for shape_num, shape in enumerate(slide.shapes, start=1):
                    if hasattr(shape, "image"):
                        img_format = Image.open(BytesIO(shape.image.blob)).format.lower()
                        img_path = os.path.join(self.temp_dir, f"slide_{slide_num}_img_{shape_num}.{img_format}")
                        with open(img_path, "wb") as img_file:
                            img_file.write(shape.image.blob)
                        image_paths.append((slide_num, img_path))
            logger.info(f"Extracted {len(image_paths)} images.")
            return image_paths
        except Exception as e:
            logger.error(f"Error extracting images from PPTX: {e}")
            raise

    @staticmethod
    def encode_image(image_path):
        try:
            with open(image_path, "rb") as img_file:
                return base64.b64encode(img_file.read()).decode("utf-8")
        except Exception as e:
            logger.error(f"Error encoding image {image_path}: {e}")
            raise

    @staticmethod
    def insert_table_into_ppt(prs, clean_rows, original_slide_num, inserted_slides_count):
        try:
            adjusted_slide_num = original_slide_num + inserted_slides_count
            new_slide = prs.slides.add_slide(prs.slide_layouts[5])  
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            xml_slides.insert(adjusted_slide_num, slides[-1])
            rows, cols = len(clean_rows), len(clean_rows[0])
            table = new_slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(8), Inches(3)).table
            for r, row in enumerate(clean_rows):
                for c, cell in enumerate(row):
                    table.cell(r, c).text = cell
            logger.info(f"Table inserted after slide {original_slide_num}.")
            return inserted_slides_count + 1
        except Exception as e:
            logger.error(f"Error inserting table into PPT: {e}")
            raise

    def process_image_with_gemini(self, image_path, slide_num):
        try:
            base64_image = self.encode_image(image_path)
            prompt = """
                You are an AI assistant that extracts handwritten tables from images.
                - Perform OCR to detect and extract table data.
                - Return the extracted table in CSV format (comma-separated values).
                - **Do not** include any Markdown syntax or separators (e.g., `|`, `---`, etc.).
                - Ensure the table structure is preserved correctly.
            """
            response = client.chat.completions.create(
                model="gemini-2.0-flash",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                        ],
                    }
                ],
            )

            extracted_text = response.choices[0].message.content.strip()
            if not extracted_text:
                logger.warning(f"No text extracted from image in slide {slide_num}.")
                return self.inserted_slides_count

            # --- Convert CSV Text to List Format ---
            clean_rows = [row.strip().split(",") for row in extracted_text.split("\n") if row.strip()]

            # --- Insert Table into PPTX ---
            self.inserted_slides_count = self.insert_table_into_ppt(self.prs, clean_rows, slide_num, self.inserted_slides_count)
            return self.inserted_slides_count
        
        except Exception as e:
            logger.error(f"Error processing image with Gemini: {e}")
            raise

    def save_pptx(self, output_pptx):
        try:
            self.prs.save(output_pptx)
            logger.info(f"Final Updated PPTX saved as {output_pptx}")
        except Exception as e:
            logger.error(f"Error saving the updated PPTX: {e}")
            raise

    def process_pptx(self):
        try:
            image_paths = self.extract_images_from_pptx()
            for slide_num, img_path in image_paths:
                self.inserted_slides_count = self.process_image_with_gemini(img_path, slide_num)
            self.save_pptx("output.pptx")
        except Exception as e:
            logger.error(f"Error during PPTX processing: {e}")
            raise

if __name__ == "__main__":
    pptx_processor = PowerPointProcessor("test.pptx")
    pptx_processor.process_pptx()