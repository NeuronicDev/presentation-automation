from fastapi import APIRouter, HTTPException, Body, status
from pydantic import BaseModel, Field
import base64, os, shutil, aiofiles, logging
from pptx import Presentation
from utils.utils import convert_pptx_to_pdf, generate_slide_context

router = APIRouter()

# --- Configuration ---
SAVE_DIR = "./uploaded_pptx"
os.makedirs(SAVE_DIR, exist_ok=True)
SLIDE_IMAGE_DIR = "./uploaded_pptx/slide_images"

# --- Pydantic Model ---
class PPTXPayload(BaseModel):
    base64: str = Field(..., description="Base64 encoded content of the PPTX file")
    filename: str = "presentation.pptx"

# --- Logging ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- Helper Functions ---
def clear_directory_contents(directory_path: str):
    """Removes all files and subdirectories within the specified directory."""
    if not os.path.isdir(directory_path):
        logger.info(f"Directory {directory_path} does not exist, nothing to clear.")
        return
    logger.info(f"Clearing contents of directory: {directory_path}")
    for filename in os.listdir(directory_path):
        file_path = os.path.join(directory_path, filename)
        try:
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
                logger.debug(f"Deleted file: {file_path}")
            elif os.path.isdir(file_path):
                shutil.rmtree(file_path)
                logger.debug(f"Deleted directory: {file_path}")
        except Exception as e:
            logger.error(f"Failed to delete {file_path}. Reason: {e}")

cleared_once = False

# --- Upload + Process Route ---
@router.post("", status_code=status.HTTP_200_OK, response_model=dict)
async def upload_pptx(payload: PPTXPayload = Body(...)):
    """
    Receives a Base64-encoded PPTX file, saves it, then processes it:
    - Converts to PDF
    - Generates slide image + XML for each slide
    - Saves image bytes in a .txt file for each slide
    """
    global cleared_once
    try:
        safe_filename = os.path.basename(payload.filename)
        if not safe_filename.lower().endswith(".pptx"):
             raise HTTPException(status_code=400, detail="Invalid file type. Only .pptx supported.")
        pptx_path = os.path.join(SAVE_DIR, safe_filename)
        pptx_bytes = base64.b64decode(payload.base64)

        async with aiofiles.open(pptx_path, "wb") as f:
            await f.write(pptx_bytes)
        logger.info(f"PPTX file saved temporarily at {pptx_path}")

        # output folder
        pptx_name = os.path.splitext(safe_filename)[0]
        slide_dir = os.path.join(SLIDE_IMAGE_DIR, pptx_name)
        pdf_output_dir = os.path.join(slide_dir, "converted_pdfs")
        os.makedirs(pdf_output_dir, exist_ok=True)

        clear_directory_contents(slide_dir)
        os.makedirs(pdf_output_dir, exist_ok=True)
        logger.info(f"Cleared and ensured directories exist for: {slide_dir}")

        # Convert to PDF
        pdf_path = convert_pptx_to_pdf(pptx_path, output_dir=pdf_output_dir)
        prs = Presentation(pptx_path)

        # Generate context for each slide
        for idx, _ in enumerate(prs.slides):
            slide_number = idx
            logger.info(f"Processing slide {slide_number}...")
            generate_slide_context(prs, slide_number, pdf_path, slide_dir)

        return {
            "status": "success",
            "message": f"File saved and processed: {safe_filename}",
            "slides_processed": len(prs.slides)
        }

    except base64.binascii.Error:
        logger.error("Base64 decoding error", exc_info=True)
        raise HTTPException(status_code=400, detail="Invalid Base64 data")
    except FileNotFoundError as e:
        logger.error("File conversion error", exc_info=True)
        raise HTTPException(status_code=500, detail=f"PDF not generated: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Server error: {e}")