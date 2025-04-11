import sys, os, logging, json, time, datetime, asyncio, pathlib, subprocess, tempfile, shutil
import logging.handlers
from pptx import Presentation
from feedback_parsing.feedback_extraction_notes import extract_feedback_from_ppt_notes
from feedback_parsing.feedback_extraction_onslide import extract_feedback_from_ppt_onslide
from feedback_parsing.feedback_extraction_mail import extract_feedback_from_email
from feedback_parsing.feedback_classifier import classify_feedback_instructions
from agents.formatting_agent import formatting_agent
from agents.cleanup_agent import cleanup_agent
from agents.visual_enhancement_agent import visual_enhancement_agent
from utils.utils import generate_slide_context, update_slide_context_cache, convert_pptx_to_pdf, extract_slide_xml_from_ppt
from utils.validation import validate_presentation
from code_manipulation.pptx_utils_functions import get_llm_mapped_function, find_shape_by_hint, function_map
from code_manipulation.code_generator import generate_python_code
from code_manipulation.code_executor import execute_code_in_docker
from code_manipulation.xml_code_generator import generate_modified_xml_code
from code_manipulation.xml_code_injector import inject_xml_into_ppt

from dotenv import load_dotenv
load_dotenv()

def setup_logging(log_file):
    logger = logging.getLogger()
    logger.setLevel(logging.INFO)
    logger.handlers.clear()
    
    file_handler = logging.handlers.RotatingFileHandler(log_file, encoding="utf-8")
    file_handler.setLevel(logging.INFO)
        
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.INFO)
    
    formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
    file_handler.setFormatter(formatter)
    console_handler.setFormatter(formatter)
    
    logger.addHandler(file_handler)
    logger.addHandler(console_handler)
    
setup_logging("presentation_automation_logs.log")


def process_with_predefined_functions(prs, task_specifications):
    remaining_tasks = []
    modified = False
    success_count = 0
    total_tasks = len(task_specifications)
    
    slide_functions = ["set_background_color"]
    table_functions = ["standardize_table", "set_table_header", "set_table_cell_color", "set_table_border"]
    
    for task in task_specifications:
        slide_number = task["slide_number"]
        if slide_number < 1 or slide_number > len(prs.slides):
            logging.warning(f"Invalid slide number: {slide_number}, skipping task")
            remaining_tasks.append(task)
            continue
        
        slide = prs.slides[slide_number - 1]
        target_hint = task.get("target_element_hint", "")
        
        mapping = get_llm_mapped_function(task)
        if mapping:
            func_name = mapping["function_name"]
            args = mapping["arguments"]
            
            if func_name in function_map:
                func = function_map[func_name]
                try:
                    if func_name in slide_functions:
                        success = func(slide, **args)
                    
                    elif func_name in table_functions:
                        success = func(slide, **args)
                    
                    else:
                        shape = find_shape_by_hint(slide, target_hint)
                        if shape:
                            if func_name == "remove_shape":
                                success = func(shape)
                            else:
                                success = func(shape, **args)
                        else:
                            logging.warning(f"No shape found for '{func_name}' on slide {slide_number}")
                            remaining_tasks.append(task)
                            continue
                    if success:
                        logging.info(f"Applied predefined function '{func_name}' to slide {slide_number}")
                        modified = True
                        success_count += 1
                    else:
                        logging.warning(f"Function '{func_name}' failed to modify slide {slide_number}")
                        remaining_tasks.append(task)
                        
                except Exception as e:
                    logging.error(f"Error applying function '{func_name}': {e}")
                    remaining_tasks.append(task)
            else:
                logging.warning(f"Function '{func_name}' not found")
                remaining_tasks.append(task)
        else:
            logging.info(f"No predefined function mapped for task: {task.get('task_description', 'Unknown')}")
            remaining_tasks.append(task)
    
    return remaining_tasks, success_count, total_tasks, modified


def process_with_python_pptx(pptx_path, base_filename_of_original_ppt, task_specifications, slide_context_cache):
    logging.info("------------------- Starting python-pptx based automation -------------------")
    
    logging.info("Generating Python code for task specifications...")
    tasks_with_code = []
    for task_specification in task_specifications:
        slide_number = task_specification["slide_number"]
        slide_context = slide_context_cache[slide_number]
        code = generate_python_code(task_specification, slide_context)
        if code:
            tasks_with_code.append({
                "slide_number": slide_number,
                "generated_code": code,
                "original_instruction": task_specification["original_instruction"],
                "description": task_specification["task_description"],
                "action": task_specification["action"],
                "target_element_hint": task_specification.get("target_element_hint", ""),
            })
            logging.info(f"Generated code for task: {task_specification['task_description']}")
        else:
            logging.warning(f"Failed to generate code for task: {task_specification['task_description']}")

    logging.info(f"Generated {len(tasks_with_code)} code snippets for the given task specifications.")


    logging.info("Executing Code in Docker...")
    execution_success, final_output_file_path, execution_report = execute_code_in_docker(tasks_with_code, pptx_path, base_filename_of_original_ppt)
    
    if execution_success:
        logging.info(f"Code execution completed. Status: {execution_report.get('status', 'unknown')}")
        if final_output_file_path:
            logging.info(f"Modified presentation saved to: {final_output_file_path}")
        else: 
            logging.error("Execution reported success but no output file path was returned.")
            execution_success = False 

        if execution_report.get("errors"):
            logging.warning(f"Some code snippets failed during execution ({len(execution_report['errors'])}):")
            for error_detail in execution_report["errors"]:
                logging.warning(f"  - Task Index {error_detail.get('task_index', 'N/A')}: {error_detail.get('error', 'Unknown error')}")
    else:
        logging.error("Code execution in Docker failed.")
        if execution_report.get("errors"):
            logging.error("Details of execution failure:")
            for error_detail in execution_report["errors"]:
                logging.error(f"  - Task Index {error_detail.get('task_index', 'N/A')}: {error_detail.get('error', 'Unknown error')}")
                        
    return execution_success, final_output_file_path, execution_report


def process_with_xml(pptx_path, base_filename_of_original_ppt, task_specifications, slide_context_cache):
    logging.info("------------------- Starting XML-based processing as fallback -------------------")
    
    execution_report = {
        "status": "unknown", 
        "errors": [], 
        "processed_count": 0, 
        "success_count": 0
    }
    
    try:
        temp_dir = tempfile.mkdtemp(prefix="ppt_xml_")
        working_pptx_path = os.path.join(temp_dir, "working_copy.pptx")
        shutil.copy2(pptx_path, working_pptx_path)
        
        slide_tasks = {}
        for task in task_specifications:
            slide_number = task["slide_number"]
            if slide_number not in slide_tasks:
                slide_tasks[slide_number] = []
            slide_tasks[slide_number].append(task)
        
        for slide_number, tasks in slide_tasks.items():
            logging.info(f"Processing {len(tasks)} tasks for slide {slide_number}")
            
            current_xml = extract_slide_xml_from_ppt(working_pptx_path, slide_number)
            if not current_xml:
                logging.error(f"Failed to extract XML for slide {slide_number}. Skipping all tasks for this slide.")
                for task in tasks:
                    execution_report["errors"].append({
                        "task_index": execution_report["processed_count"],
                        "error": f"Failed to extract XML for slide {slide_number}",
                        "description": task.get("task_description", "Unknown")
                    })
                    execution_report["processed_count"] += 1
                continue
                
            slide_context = slide_context_cache[slide_number]
            
            # Apply each task sequentially
            for i, task in enumerate(tasks):
                task_desc = task.get("task_description", f"Task #{i+1}")
                logging.info(f"Applying task {i+1}/{len(tasks)} for slide {slide_number}: {task_desc}")
                
                modified_xml = generate_modified_xml_code(
                    original_xml=current_xml,
                    agent_task_specification=task,
                    slide_context=slide_context
                )
                
                execution_report["processed_count"] += 1
                
                if not modified_xml:
                    logging.error(f"Failed to generate modified XML for task: {task_desc}")
                    execution_report["errors"].append({
                        "task_index": execution_report["processed_count"] - 1,
                        "error": "Failed to generate modified XML",
                        "description": task_desc
                    })
                    continue
                    
                if modified_xml == current_xml:
                    logging.warning(f"No XML changes made for task: {task_desc}")
                    execution_report["success_count"] += 1
                    continue
                
                current_xml = modified_xml
                
                success = inject_xml_into_ppt(working_pptx_path, slide_number, modified_xml)
                
                if success:
                    logging.info(f"Successfully applied XML changes for task: {task_desc}")
                    execution_report["success_count"] += 1
                else:
                    logging.error(f"Failed to inject modified XML for task: {task_desc}")
                    execution_report["errors"].append({
                        "task_index": execution_report["processed_count"] - 1,
                        "error": "Failed to inject modified XML",
                        "description": task_desc
                    })
        
        final_output_filename = f"{base_filename_of_original_ppt}_xml_modified.pptx"
        output_dir = os.path.abspath("./output_ppts/pptx")
        os.makedirs(output_dir, exist_ok=True)
        final_output_path = os.path.join(output_dir, final_output_filename)
        
        shutil.copy2(working_pptx_path, final_output_path)
        logging.info(f"Final modified presentation saved to: {final_output_path}")

        shutil.rmtree(temp_dir)
        
        if execution_report["success_count"] == execution_report["processed_count"]:
            execution_report["status"] = "success"
        elif execution_report["success_count"] > 0:
            execution_report["status"] = "partial_success"
        else:
            execution_report["status"] = "failed"
            
        execution_success = execution_report["success_count"] > 0
        return execution_success, final_output_path, execution_report
        
    except Exception as e:
        logging.error(f"XML processing failed with unexpected error: {e}", exc_info=True)
        execution_report["status"] = "failed"
        execution_report["errors"].append({
            "task_index": -1, 
            "error": f"Global XML processing error: {str(e)}",
            "description": "Overall XML processing"
        })
        return False, None, execution_report


def main(original_pptx_path, email_path="path/to/email.txt"):
    overall_pipeline_status = "unknown"
    
    temp_dir = None
    try:
        logging.info("======================= Starting PPT Automation Pipeline ========================")
        logging.info(f"Input PPTX: {original_pptx_path}")
        base_filename_of_original_ppt = os.path.splitext(os.path.basename(original_pptx_path))[0]
        
        # Create temporary directory for intermediate files
        temp_dir = tempfile.mkdtemp(prefix="ppt_pipeline_")
        logging.info(f"Created temporary directory: {temp_dir}")
        temp_pdf_dir = os.path.join(temp_dir, "pdf")
        os.makedirs(temp_pdf_dir, exist_ok=True)
        
        # pdf_path = convert_pptx_to_pdf(pptx_path, output_dir = os.path.abspath("./input_ppts/converted_pdfs"))
        pdf_path = convert_pptx_to_pdf(original_pptx_path, output_dir=temp_pdf_dir)  
        
        logging.info("Loading presentation...")
        prs = Presentation(original_pptx_path)

        image_cache = {}
        slide_context_cache = {}

        # Step 1: Feedback Extraction
        logging.info("Extracting feedback from all sources...")
        feedback_notes = extract_feedback_from_ppt_notes(original_pptx_path)
        feedback_onslide = extract_feedback_from_ppt_onslide(original_pptx_path)
        feedback_mail = extract_feedback_from_email(email_path)
        
        all_feedback = feedback_notes + feedback_onslide + feedback_mail
        logging.info(f"Total feedback instructions extracted: {len(all_feedback)}")

        # Step 2: Instruction Interpretation
        logging.info("Classifying and extracting tasks from feedback...")
        categorized_tasks = classify_feedback_instructions(all_feedback)
        logging.info(f"Categorized Tasks: {len(categorized_tasks)}")

        # Step 3: Delegate tasks to specialized agents
        logging.info("Processing Tasks with Agents & Context...")
        task_specifications = []
        for task in categorized_tasks:
            category = task["category"]
            slide_number = task["slide_number"]
            instruction = task["original_instruction"]
            
            # Generate or retrieve slide context
            if slide_number not in slide_context_cache:
                slide_context_cache[slide_number] = generate_slide_context(prs, slide_number, pdf_path, image_cache)
            slide_context = slide_context_cache[slide_number]
            
            # Delegate to appropriate agent
            if category == "formatting":
                task_with_desc = formatting_agent(task, slide_context)
            elif category == "cleanup":
                task_with_desc = cleanup_agent(task, slide_context)
            elif category == "visual_enhancement":
                task_with_desc = visual_enhancement_agent(task, slide_context)
            else:
                logging.warning(f"Unknown category: {category} for instruction: '{instruction}'")
                continue
            
            if task_with_desc:
                task_specifications.extend(task_with_desc)
    
        if not task_specifications:
            logging.warning("No tasks generated from the feedback. Nothing to process.")
            return None, {"status": "no_tasks", "message": "No tasks to process"}
        
        
        predefined_execution_report = {
            "status": "unknown",
            "processed_count": len(task_specifications),
            "success_count": 0,
            "errors": []
        }
                
        # Step 4: Process with predefined functions first
        logging.info("Attempting to process tasks with predefined functions...")
        remaining_tasks, predefined_success_count, total_tasks, modified = process_with_predefined_functions(prs, task_specifications)
        predefined_execution_report["success_count"] = predefined_success_count
        
        intermediate_path = None
        if modified:
            # Save intermediate changes to temp directory
            intermediate_filename = f"intermediate_predefined.pptx"
            intermediate_path = os.path.join(temp_dir, intermediate_filename)
            prs.save(intermediate_path)
            logging.info(f"Intermediate presentation with predefined changes saved to: {intermediate_path}")
            
            if predefined_success_count == total_tasks:
                predefined_execution_report["status"] = "success"
            elif predefined_success_count > 0:
                predefined_execution_report["status"] = "partial_success"
            else:
                predefined_execution_report["status"] = "failed"

            if not remaining_tasks:
                logging.info("All tasks were successfully processed with predefined functions.")
                
                # Copy to permanent output location
                final_output_filename = f"{base_filename_of_original_ppt}_predefined_func_modified.pptx"
                output_dir = os.path.abspath("./output_ppts/pptx")
                os.makedirs(output_dir, exist_ok=True)
                final_output_path = os.path.join(output_dir, final_output_filename)
                shutil.copy2(intermediate_path, final_output_path)
                
                validation_report, validation_success, validation_success_percentage = validate_presentation(
                    original_pptx_path,
                    final_output_path,
                    task_specifications
                )
                if validation_success_percentage < 60:
                    logging.warning(f"Validation success rate too low ({validation_success_percentage:.1f}%). Falling back to pptx approach...")
                else: 
                    if not validation_success:
                        logging.warning(f"Validation failed, success percentage is {validation_success_percentage:.1f}%. Some modifications may not have been applied correctly.")
                        overall_pipeline_status = "partial_success"
                    else:
                        overall_pipeline_status = "success"
                    
                    final_report = {
                        "approach": "predefined_functions",
                        "execution_report": predefined_execution_report,
                        "validation_report": validation_report,
                        "final_status": overall_pipeline_status,
                        "success_rate": (predefined_success_count / total_tasks) * 100,
                        "output_path": final_output_path
                    }
                    return final_output_path, final_report
            else:
                logging.info(f"Some {predefined_success_count} tasks were processed with predefined functions out of {total_tasks} , but {len(remaining_tasks)} tasks remain.")
                slide_context_cache, pdf_path = update_slide_context_cache(intermediate_path, pdf_path, slide_context_cache, image_cache)

        # Update pptx_path for further processing
        current_pptx_path = intermediate_path if intermediate_path else original_pptx_path

                
        if remaining_tasks:
            # Step 5: Fall back to python-pptx for remaining tasks
            if predefined_success_count > 0:
                logging.info(f"Processed {predefined_success_count}/{total_tasks} tasks with predefined functions. {len(remaining_tasks)} tasks remain. Falling back to python-pptx approach...")
            else:
                logging.info(f"All {total_tasks} tasks were not processed with predefined functions. {len(remaining_tasks)} tasks remain. Falling back to python-pptx approach...")
            
            logging.info("Attempting to process tasks with python-pptx approach...")            
            pptx_success, pptx_output_path, pptx_report = process_with_python_pptx(
                current_pptx_path, 
                base_filename_of_original_ppt,
                remaining_tasks, 
                slide_context_cache
            )
            
            if pptx_success and pptx_report.get("status") in ["success", "partial_success"]:
                logging.info("Python-pptx approach was successful. Using its output.")
                
                success_rate = pptx_report.get("success_count", 0) / len(remaining_tasks) * 100
                
                if pptx_report["status"] == "partial_success" and success_rate < 60:
                    logging.info(f"Success rate of python-pptx approach was only {success_rate:.1f}%. Trying XML approach for a possibly better result.")
                else:
                    final_output_path = pptx_output_path
                    
                    logging.info("Validating modifications...")
                    validation_report, validation_success, validation_success_percentage = validate_presentation(
                        original_pptx_path,
                        final_output_path,
                        task_specifications
                    )
                    
                    if validation_success_percentage < 60:
                        logging.warning(f"Validation success rate too low ({validation_success_percentage:.1f}%). Falling back to XML approach...")

                    else: 
                        if not validation_success:
                            logging.warning(f"Validation failed, success percentage is {validation_success_percentage:.1f}% . Some modifications may not have been applied correctly.")
                            overall_pipeline_status = "partial_success"
                        else:
                            overall_pipeline_status = pptx_report["status"]
                        
                        final_report = {
                            "approach": "python-pptx",
                            "predefined_execution_report": predefined_execution_report,
                            "pptx_execution_report": pptx_report,
                            "final_status": overall_pipeline_status,
                            "success_rate": success_rate,
                            "validation_report": validation_report,
                            "output_path": final_output_path,
                        }
                        return final_output_path, final_report
            
            # Step 6: Fallback to XML approach if python-pptx failed or had low success rate
            logging.info("Falling back to XML approach...")
            xml_success, xml_output_path, xml_report = process_with_xml(
                current_pptx_path, 
                base_filename_of_original_ppt,
                remaining_tasks, 
                slide_context_cache
            )
            
            # Determine which result to use
            if not pptx_success and not xml_success:
                logging.error("Both approaches failed. No successful output generated.")
                final_output_path = None
                overall_pipeline_status = "failed"
                final_report = {
                    "approach": "both_failed",
                    "pptx_report": pptx_report,
                    "xml_report": xml_report,
                    "final_status": "failed"
                }
                
            elif not pptx_success and xml_success:
                logging.info("XML approach succeeded where python-pptx failed. Using XML output.")
                final_output_path = xml_output_path
                
                logging.info("Validating XML modifications...")
                validation_report, validation_success, validation_success_percentage = validate_presentation(
                    original_pptx_path,
                    final_output_path,
                    task_specifications
                )
                
                if not validation_success:
                    logging.warning("Validation failed. Some modifications may not have been applied correctly.")
                    overall_pipeline_status = "partial_success"
                else: 
                    overall_pipeline_status = xml_report["status"]
                
                final_report = {
                    "approach": "xml",
                    "predefined_execution_report": predefined_execution_report,
                    "xml_execution_report": xml_report,
                    "final_status": overall_pipeline_status,
                    "success_rate": xml_report.get("success_count", 0) / len(remaining_tasks) * 100,
                    "validation_report": validation_report,
                    "output_path": final_output_path,
                }
                
            else:
                # Both succeeded or python-pptx partially succeeded - compare results
                pptx_success_rate = pptx_report.get("success_count", 0) / len(remaining_tasks) * 100
                xml_success_rate = xml_report.get("success_count", 0) / len(remaining_tasks) * 100
                
                if xml_success_rate >= pptx_success_rate :
                    logging.info(f"XML approach had better success rate ({xml_success_rate:.1f}% vs {pptx_success_rate:.1f}%). Using XML output.")
                    final_output_path = xml_output_path
                    
                    logging.info("Validating XML modifications...")
                    validation_report, validation_success, validation_success_percentage = validate_presentation(
                        original_pptx_path,
                        final_output_path,
                        task_specifications
                    )
                    
                    if not validation_success:
                        logging.warning("Validation failed. Some modifications may not have been applied correctly.")
                        overall_pipeline_status = "partial_success"
                    else:
                        overall_pipeline_status = xml_report["status"]
                        
                    final_report = {
                        "approach": "xml_better",
                        "predefined_execution_report": predefined_execution_report,
                        "pptx_report": pptx_report,
                        "xml_report": xml_report,
                        "final_status": overall_pipeline_status,
                        "success_rate": xml_success_rate,
                        "validation_report": validation_report,
                        "output_path": xml_output_path,
                    }
                else: 
                    logging.info(f"Python-pptx approach had better success rate ({pptx_success_rate:.1f}% vs {xml_success_rate:.1f}%). Using python-pptx output.")
                    final_output_path = pptx_output_path
                    
                    logging.info("Validating XML modifications...")
                    validation_report, validation_success, validation_success_percentage = validate_presentation(
                        original_pptx_path,
                        final_output_path,
                        task_specifications
                    )
                    
                    if not validation_success:
                        logging.warning("Validation failed. Some modifications may not have been applied correctly.")
                        overall_pipeline_status = "partial_success"
                    else:
                        overall_pipeline_status = pptx_report["status"]

                    final_report = {
                        "approach": "pptx_better",
                        "predefined_execution_report": predefined_execution_report,
                        "pptx_report": pptx_report,
                        "xml_report": xml_report,
                        "final_status": overall_pipeline_status,
                        "success_rate": pptx_success_rate,
                        "validation_report": validation_report,
                        "output_path": pptx_output_path,
                    }
            
            return final_output_path, final_report

    except FileNotFoundError as fnf_err:
        logging.error(f"File not found error during pipeline: {fnf_err}")
        overall_pipeline_status = "failed"
        return None, {"status": "failed", "error": str(fnf_err)}
    except Exception as e:
        logging.error(f"Pipeline failed with an unexpected error: {e}", exc_info=True)
        overall_pipeline_status = "failed"
        return None, {"status": "failed", "error": str(e)}
    finally:
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
                logging.info(f"Cleaned up temporary directory: {temp_dir}")
            except Exception as e:
                logging.warning(f"Failed to clean up temporary directory: {e}")
        
        logging.info(f"================== Pipeline Finished with Overall Status: {overall_pipeline_status} ===================")


if __name__ == "__main__":
    # original_pptx_path = os.path.abspath("./input_ppts/pptx/font_test_pfunc.pptx")
    # original_pptx_path = os.path.abspath("./input_ppts/pptx/font_test1.pptx")
    original_pptx_path = os.path.abspath("./input_ppts/pptx/font_test2.pptx")
    # original_pptx_path = os.path.abspath("./input_ppts/pptx/cleanup_test.pptx")
    # original_pptx_path = os.path.abspath("./input_ppts/pptx/cleanup_test2.pptx")
    # original_pptx_path = os.path.abspath("./input_ppts/pptx/table_alignment_test.pptx")
    # original_pptx_path = os.path.abspath("./input_ppts/pptx/table_alignment_test2.pptx")
    # original_pptx_path = os.path.abspath("./input_ppts/pptx/consistent.pptx")
    # original_pptx_path = os.path.abspath("./input_ppts/pptx/consistent3.pptx")

    logging.info(f"Input PPTX: {original_pptx_path}")
    output_path, report = main(original_pptx_path)
    
    if output_path:
        logging.info(f"Process completed. Final output at: {output_path}")
        logging.info(f"Final report: {json.dumps(report, indent=2)}")
    else:
        logging.error("Process failed. No output generated.")
    