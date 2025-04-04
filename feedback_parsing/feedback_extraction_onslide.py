import logging
import os
from typing import List, Dict, Optional, Union, Tuple
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.dml.color import RGBColor
from pptx.exc import PackageNotFoundError

DEFAULT_ANNOTATION_KEYWORDS = ["feedback:", "note:", "todo:", "comment:"]
DEFAULT_ANNOTATION_COLORS = [
    RGBColor(255, 255, 0),  # Yellow
    RGBColor(255, 192, 0),   # Amber/Orange-Yellow
    RGBColor(255, 0, 0),     # Red
    RGBColor(0, 255, 0),     # Green
]

def _is_annotation_color(shape_fill, target_colors: List[RGBColor]) -> bool:
    if shape_fill.type == MSO_FILL_TYPE.SOLID:
        try:
            if hasattr(shape_fill, 'fore_color') and hasattr(shape_fill.fore_color, 'rgb'):
                if shape_fill.fore_color.rgb is None:
                    return False
                return shape_fill.fore_color.rgb in target_colors
        except AttributeError:
            logging.debug("AttributeError checking shape fill color.")
            return False
    return False


def has_text(shape) -> bool:
    try:
        return shape.has_text_frame and shape.text_frame.text.strip() != ""
    except Exception:
        return False

def extract_text(shape) -> str:
    try:
        return shape.text_frame.text.strip()
    except Exception:
        return ""
    

def get_position_details(shape) -> Dict:
    try:
        return {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height
        }
    except AttributeError:
        return {}


def extract_feedback_from_ppt_onslide(pptx_filepath: str, annotation_keywords: Optional[List[str]] = None, annotation_colors: Optional[List[RGBColor]] = None ) -> List[Dict[str, Union[str, int, Dict]]]:
    logging.info(f"Extracting feedback from sticky notes")
    feedback_list = []
    # keywords = annotation_keywords if annotation_keywords is not None else DEFAULT_ANNOTATION_KEYWORDS
    # colors = annotation_colors if annotation_colors is not None else DEFAULT_ANNOTATION_COLORS

    # try:
    #     prs = Presentation(pptx_filepath)
    #     for i, slide in enumerate(prs.slides):
    #         slide_number = i + 1
    #         logging.info(f"Processing slide {slide_number}...")
            
    #         for shape in slide.shapes:
    #             is_feedback = False
    #             feedback_text = ""

    #             if has_text(shape):
    #                 shape_text = extract_text(shape)
    #                 shape_text_lower = shape_text.lower()
                    
    #                 for keyword in keywords:
    #                     if shape_text_lower.startswith(keyword):
    #                         feedback_text = shape_text
    #                         is_feedback = True
    #                         logging.info(f"Slide {slide_number}: Found feedback by keyword '{keyword}' (Shape ID: {shape.shape_id})")
    #                         break
                    
    #                 if not is_feedback and hasattr(shape, 'fill'):
    #                     if _is_annotation_color(shape.fill, colors) and shape_text.lower() in keywords:
    #                         feedback_text = shape_text
    #                         is_feedback = True
    #                         logging.info(f"Slide {slide_number}: Found feedback by color (Shape ID: {shape.shape_id})")
                
    #             if is_feedback and feedback_text:
    #                 pos_details = get_position_details(shape)
                    
    #                 feedback_item = {
    #                     "source": "on_slide",
    #                     "slide_number": slide_number,
    #                     "feedback": feedback_text,
    #                     "element_details": {
    #                         "shape_id": getattr(shape, 'shape_id', None),
    #                         "shape_name": getattr(shape, 'name', ""), 
    #                         "text": feedback_text,
    #                         "position": pos_details
    #                     }
    #                 }
    #                 feedback_list.append(feedback_item)
            
    #         logging.info(f"Finished processing slide {slide_number}.")
        
    #     logging.info(f"Extraction complete. Found {len(feedback_list)} feedback items in onslide sticky notes.")
        
    # except PackageNotFoundError:
    #     logging.error(f"Error: Not a valid PPTX file: {pptx_filepath}")
    # except Exception as e:
        # logging.error(f"An unexpected error occurred during extraction: {e}", exc_info=True)
    
    return feedback_list


# if __name__ == "__main__":
#     logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
#     pptx_filepath = "Source.pptx"
#     feedback = extract_feedback_from_ppt_onslide(pptx_filepath)
#     print(feedback)
