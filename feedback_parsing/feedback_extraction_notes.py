import logging
from typing import List, Dict, Optional, Union
from pptx import Presentation
from pptx.exc import PackageNotFoundError

def extract_feedback_from_ppt_notes(pptx_filepath: str) -> List[Dict[str, Union[str, int, None]]]:
    feedback_list: List[Dict[str, Union[str, int, None]]] = []
    logging.info(f"Starting feedback extraction from notes in: {pptx_filepath}")

    try:
        prs = Presentation(pptx_filepath)
        logging.info(f"Successfully opened presentation. Found {len(prs.slides)} slides.")

        for i, slide in enumerate(prs.slides):
            slide_number = i + 1  
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                    instruction_text = notes_slide.notes_text_frame.text.strip()
                    if instruction_text:  
                        feedback_item = {
                            "source": "notes",
                            "slide_number": slide_number,
                            "instruction": instruction_text
                        }
                        feedback_list.append(feedback_item)
                        logging.info(f"Found notes feedback on slide {slide_number}: '{instruction_text[:50]}...'")
            else:
                logging.info(f"Slide {slide_number} has no notes slide or notes text.")
                
        logging.info(f"Finished notes extraction. Found {len(feedback_list)} feedback items.")
    except PackageNotFoundError:
        logging.error(f"Error: File not found or not a valid PPTX file: {pptx_filepath}")
    except Exception as e:
        logging.error(f"An unexpected error occurred during notes extraction in extract_feedback_from_ppt_notes: {e}", exc_info=True)
        
    return feedback_list

