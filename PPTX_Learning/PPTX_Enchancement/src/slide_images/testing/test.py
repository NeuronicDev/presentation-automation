from pptx import Presentation
from pptx.util import Inches
import os

def adjust_canvas_and_save_pdf(input_path, output_path):
    prs = Presentation(input_path)
    new_prs = Presentation()

    for slide in prs.slides:
        max_x, max_y = prs.slide_width, prs.slide_height
        for shape in slide.shapes:
            if shape.shape_type != 6:  
                right = shape.left + shape.width
                bottom = shape.top + shape.height
                max_x = max(max_x, right)
                max_y = max(max_y, bottom)

        # Set new slide size for capturing overflow
        new_prs.slide_width = max_x
        new_prs.slide_height = max_y

        new_slide = new_prs.slides.add_slide(new_prs.slide_layouts[6])
        for shape in slide.shapes:
            el = shape.element
            new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

    temp_pptx = "temp_overflow.pptx"
    new_prs.save(temp_pptx)

    os.system(f'soffice --headless --convert-to pdf "{temp_pptx}" --outdir "{output_path}"')

    print(f"PDF exported to: {os.path.join(output_path, temp_pptx.replace('.pptx', '.pdf'))}")
adjust_canvas_and_save_pdf("presentation.pptx", "./")