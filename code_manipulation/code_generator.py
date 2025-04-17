import json, logging, re
from typing import Dict, Any, Optional

from langchain.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import HumanMessage

from config.llm_provider import gemini_flash_llm  
from config.config import LLM_API_KEY

from google import genai
client = genai.Client(api_key=LLM_API_KEY)

CODE_GENERATION_PROMPT = """
    # PowerPoint Automation Expert

    You are a **State-of-the-Art expert Python code generator** and PowerPoint automation engineer with deep expertise in producing error-free, production-ready code using python-pptx-1.0.0 library and PowerPoint's object model to manipulate PowerPoint presentations.
    Your task is to generate precise, executable Python code to modify a PowerPoint slide based on the provided feedback instruction, task details and slide context.
    Your generated code will be executed within a specific environment where `prs` (the Presentation object) and `slide` (the target Slide object, if applicable) are already defined in the execution scope. Use them directly.
    Ensure the generated code accurately performs the specified action on the target element(s) while preserving the original slide content.


    ## Instructions:
        - Understand the task description and original instruction provided by the agent.
        - Analyze the slide's current state using the provided image and XML structure and determine the changes to be made.
        - Identify the target element(s) based on the `target_element_hint` and the slide's context(visual image and XML structure).
        - You should not assume anything which is not in the slide context. Use image and XML structure to understand the slide and give code accordingly.
        - The MOST IMPORTANT requirement is to generate Python code that is **directly executable** and **free of errors**. Pay extreme attention to `python-pptx` syntax, correct attribute and method names, and proper import statements.
    
    ## Input Context:
        - Agent: {agent_name}
        - Slide Index: {slide_index}
        - Original Instruction: {original_instruction}
        - Task Description: {task_description}
        - Action: {action}
        - Target Element Hint: {target_element_hint}
        - Parameters: {params}
        - Slide XML Structure: {slide_xml_structure}


    ## IMPORTANT: Python-PPTX Technical Guidance & Best Practices:

    1. **Slide Properties:**
    - Access slide dimensions via `prs.slide_width` and `prs.slide_height` (NOT slide.slide_width)
    - All position/size values are in EMU units (English Metric Units)
    - Convert units with `pptx.util.Inches()`, `pptx.util.Pt()`, or `pptx.util.Cm()`

    2. **Shape Manipulation:**
    - Always verify shape properties before access: `if shape.has_text_frame:`, `if shape.has_table:`, etc.
    - Access text with shape.text_frame and paragraphs[0].runs[0].text
    - Shapes have properties: `.left`, `.top`, `.width`, `.height`, `.rotation`, `.name`, `.shape_type`, `.add_shape()` etc.
    - Shape types include: `shape.shape_type`, MSO_SHAPE_TYPE enum from pptx.enum.shapes
    - Access shapes by index: `slide.shapes[0]` or iterate: `for shape in slide.shapes:`
    - Find shapes by name: `[s for s in slide.shapes if s.name == "target_name"]`
    - Different shape types (`shape.shape_type`) have different properties.
    - Basic shapes (AutoShapes, TextBoxes) often have direct `.fill` and `.line` attributes.
    - Container shapes like `GraphicFrame` (MSO_SHAPE_TYPE.GRAPHIC_FRAME, which holds Tables, Charts, SmartArt) DO NOT have a direct `.fill`. You must access the object inside (e.g., `shape.table`, `shape.chart`) and format its components (cells, plot area).
    - **ALWAYS check `shape.shape_type` or use `hasattr(shape, 'fill')` before attempting to access attributes like `.fill` or `.line` to avoid AttributeErrors.**
    - There will be no external images/icons you should add while creating the code. You should only use the images/icons that are already present in the slide or create new basic shapes/icons internally.

    3. **Text Operations:**
    - Text frame access: `shape.text_frame.text = "New text"`
    - Paragraph formatting: `paragraph.alignment`, `paragraph.level`, `paragraph.space_before`, from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
    - Run properties: `run.font.size`, `run.font.bold`, `run.font.color.rgb`
    - Auto-size: `text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT`

    4. **Table Operations:**
    - Create tables with shapes.add_table(rows, cols, x, y, cx, cy)
    - Table access: `table = shape.table`
    - Cell access: `cell = table.cell(row_idx, col_idx)`
    - Format cells with table.cell(row, col).fill.solid(), table.cell(row, col).text_frame
    - Cell properties: `cell.text`, `cell.fill`, `cell.margin_x`
    - Table dimensions: `table.rows`, `table.columns`
    - Set line properties with shape.line.width, shape.line.color, shape.line.dash_style

    5. **Color Operations:**
    - from pptx.dml.color import RGBColor
    - RGB colors: `RGBColor(r, g, b)` where r,g,b are 0-255
    - Theme colors: `MSO_THEME_COLOR.ACCENT_1`
    - Shape fill: `shape.fill.solid()`, `shape.fill.fore_color.rgb = RGBColor(r,g,b)`
    - When writing code directly use rgb values instead of theme colors/ hex codes to avoid errors.
    
    6. **Positioning Relative to Table Cells:**
    - Individual table cells (`table.cell(r, c)`) DO NOT have `.left` or `.top` attributes.
    - To position a shape (e.g., an icon) visually within or near a cell at `(row, col)`, you MUST calculate its position relative to the TABLE's shape:
     1. Get the table shape's position: `table_shape = ...`, `table_left = table_shape.left`, `table_top = table_shape.top`.
     2. Calculate the sum of widths of columns *before* `col`: `left_offset = sum(table.columns[c].width for c in range(col))`.
     3. Calculate the sum of heights of rows *before* `row`: `top_offset = sum(table.rows[r].height for r in range(row))`.
     4. The approximate cell position is `(table_left + left_offset, table_top + top_offset)`.
     5. Add the new shape to the `slide.shapes` collection using these calculated coordinates (plus any desired padding).
   - Use `try...except` when accessing column widths and row heights as indices might be invalid. Handle potential `AttributeError` if `.width` or `.height` aren't available directly.
    
    7. **Error Prevention:**
    - ALWAYS use try/except blocks for operations that might fail
    - Check if shapes exist before modification
    - Verify attribute existence before access. Do NOT use attributes that may not exist.
    - Use relative positioning rather than hardcoded coordinates
    - Preserve existing content unless explicitly told to modify

    8. **VERIFIED IMPORTS & ENUMS - USE ONLY THESE:**    
        from pptx.util import Inches, Pt, Emu, Cm
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT, MSO_TEXT_UNDERLINE_TYPE
        from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE, MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE, PP_MEDIA_TYPE
        from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE, MSO_COLOR_TYPE, MSO_PATTERN_TYPE, MSO_THEME_COLOR_INDEX
        from pptx.chart.data import CategoryChartData, ChartData, XyChartData, BubbleChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE, XL_DATA_LABEL_POSITION
        from pptx.enum.action import PP_ACTION_TYPE
        from pptx.table import Table, _Cell 
    

    **Common Tasks:**
        *   **Text:** Access `shape.text_frame`, `tf.paragraphs`, `p.runs`, `run.font` (use `run.font.name`, `.size = Pt(...)`, `.bold = True`, `.color.rgb = RGBColor(...)`), `p.alignment = PP_ALIGN...`. Check `shape.has_text_frame` first.
        *   **Position/Size:** Use `shape.left`, `shape.top`, `shape.width`, `shape.height`. Use `Inches()` or `Emu()` for setting values. Calculate relative positions where possible.
        *   **Fill/Line:** Access `shape.fill` (`fill.solid()`, `fill.fore_color.rgb = RGBColor(...)`), `shape.line` (`line.color.rgb = ...`, `line.width = Pt(...)`). Check fill/line type first if needed (`fill.type == MSO_FILL_TYPE...`).
        *   **Adding Shapes:** Use `slide.shapes.add_textbox(...)`, `add_picture(...)`, `add_shape(MSO_AUTO_SHAPE_TYPE..., left, top, width, height)`, `add_table(...)`, `add_connector(MSO_CONNECTOR_TYPE..., begin_x, begin_y, end_x, end_y)`.
        *   **Deleting Shapes:** This requires accessing the internal XML element: `sp = shape._sp`, then `sp.getparent().remove(sp)`. Use this pattern carefully within a loop (iterate over a *copy* of shapes list if deleting).
        *   **Tables:** Access `shape.table`, `table.cell(r, c)`, `cell.text_frame`. Formatting involves iterating cells/rows/columns. Merging is `cell.merge(...)`.
        *   **Charts:** Create charts with slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data), Access `shape.chart`, `chart.series[0].values`, `chart.plots[0].categories`, `chart.plots[0].has_data_labels = True`, `chart.plots[0].data_labels.show_value = True`.  
        *   **Connectors:** After adding, use `connector.begin_connect(shape, connection_site_index)` and `connector.end_connect(...)`.
        *   **Layout:** `python-pptx` has no high-level "align top" or "distribute". You MUST calculate the required `left`, `top`, `width`, `height` values based on `prs.slide_width/height`, other shapes' positions, and the `task_description`. Example: To center horizontally: `shape.left = (prs.slide_width - shape.width) // 2`.

    **Specialized Layout Techniques:**
    - **Smart Layouts:** For timelines, use `slide.shapes.add_connector()` with calculated spacing; for matrices, standardize cell sizes.    
    - **Process Flows**: Ensure consistent spacing and sizing of each step. Connect with MSO_CONNECTOR_TYPE.STRAIGHT.
    - **Tables & Matrices**: Standardize cell sizes and ensure consistent borders using line properties.
    - **Comparison Layouts**: Maintain visual balance between compared elements with equal sizing.
    
    
    ## CRITICAL REQUIREMENTS:
    - ONLY use documented public methods and attributes from python-pptx (NEVER use internal methods that start with underscore "_" except for deletion operations where specifically mentioned)
    - NEVER invent or use non-existent methods, attributes, or enums 
    - ALWAYS use proper data types (integers for indices, strings for text, etc.)
    - ALWAYS verify your imports are correct and reflect actual python-pptx modules
    - ENSURE all enum values used are valid members of their respective enum classes
    
    ## COMMON MISTAKES TO AVOID:
    - DO NOT use non-existent methods like `shape.get_text()` or `slide.align_shapes()`
    - DO NOT use internal methods that start with underscore (e.g., `._element`) except for deletion
    - DO NOT invent new enums or use incorrect enum values
    - DO NOT use theme colors without verifying they exist; prefer RGBColor values
    - DO NOT mix up coordinate units (use Inches/Pt/Cm consistently)
    - DO NOT use string indices for accessing elements in collections
    - DO NOT assume table cells have direct position attributes
    - DO NOT forget proper type conversions (e.g., `int()` for indices)
    - DO NOT use variables before defining them

    
    ## OUTPUT FORMAT:
    Provide ONLY executable Python code without explanations, preamble, or markdown formatting. JUST the code.
    The code should be ready for immediate execution.
    """


def generate_python_code(agent_task_specification: Dict[str, Any], slide_context: Dict[str, Any]) -> Optional[str]:
    try:
        agent_name = agent_task_specification.get("agent_name")
        slide_number = agent_task_specification.get("slide_number")
        original_instruction = agent_task_specification.get("original_instruction")
        task_description = agent_task_specification.get("task_description")
        action = agent_task_specification.get("action")
        target_element_hint = agent_task_specification.get("target_element_hint")
        params = agent_task_specification.get("params", {})
        slide_xml_structure = slide_context.get("slide_xml_structure")
        slide_image_base64 = slide_context.get("slide_image_base64") 
        slide_image_bytes = slide_context.get("slide_image_bytes")  

        main_prompt = CODE_GENERATION_PROMPT.format(
            agent_name=agent_name,
            slide_index=slide_number - 1,
            original_instruction=original_instruction,
            task_description=task_description,
            action=action,
            target_element_hint=target_element_hint,
            params=params,
            slide_xml_structure=slide_xml_structure,
        )       
        
        final_prompt = [main_prompt]
        slide_image_text_prompt = "The below is visual representation of the slide image. Use it to understand and visualize the current layout, structure, elements, spacing, overlaps, colors, and styles."
        final_prompt.append(slide_image_text_prompt)

        image = genai.types.Part.from_bytes(data=slide_image_bytes, mime_type="image/png") 
                
        try:
            response = client.models.generate_content(model="gemini-2.0-flash", contents=[final_prompt, image])
            generated_code_str = response.text.strip()
            code_block = re.search(r'```python\n(.*?)\n```', generated_code_str, re.DOTALL)
            
            if code_block:
                extracted_code = code_block.group(1)
                logging.info(f"GENERATED CODE:\n {extracted_code}")
                return extracted_code
            else:
                logging.info(f"No markdown code block found. Using entire string.")
                return generated_code_str
        except Exception as e:
            logging.info(f"Error generated code: {e}")
            raise 
    
    except Exception as e:
        logging.error(f"Error generating Python code for task: {agent_task_specification.get('action')}. {e}", exc_info=True)
        return None

