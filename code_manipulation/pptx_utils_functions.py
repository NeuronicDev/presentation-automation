import json, logging, re, os, sys, time, datetime, asyncio, subprocess, tempfile, shutil, pathlib
import logging.handlers
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE
from pptx.oxml.xmlchemy import OxmlElement

from config.config import LLM_API_KEY
from google import genai
client = genai.Client(api_key=LLM_API_KEY)

# ----- TEXT FORMATTING FUNCTIONS -----
def set_font(shape, font_name="Arial", font_size=12):
    """Change the font name and size of text in a shape."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                if font_size is not None:
                    run.font.size = Pt(font_size)
        return True
    return False

def align_text(shape, alignment="center"):
    """Align text in a shape (left, center, right)."""
    if hasattr(shape, 'text_frame'):
        alignment_map = {"left": PP_PARAGRAPH_ALIGNMENT.LEFT, "center": PP_PARAGRAPH_ALIGNMENT.CENTER, "right": PP_PARAGRAPH_ALIGNMENT.RIGHT}
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = alignment_map.get(alignment.lower(), PP_PARAGRAPH_ALIGNMENT.CENTER)
        return True
    return False

def adjust_spacing(shape, line_spacing=1.0):
    """Adjust line spacing in a shape's text frame."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.line_spacing = line_spacing
        return True
    return False

def set_text_color(shape, color_rgb=(0, 0, 0)):
    """Set text color in a shape using RGB values."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if isinstance(color_rgb, list):
                    color_rgb = tuple(color_rgb)
                run.font.color.rgb = RGBColor(*color_rgb)
        return True
    return False

def apply_bold(shape):
    """Apply bold formatting to text in a shape."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
        return True
    return False

def apply_italic(shape):
    """Apply italic formatting to text in a shape."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.italic = True
        return True
    return False

def remove_bold(shape):
    """Remove bold formatting from text in a shape."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = False
        return True
    return False

def remove_italic(shape):
    """Remove italic formatting from text in a shape."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.italic = False
        return True
    return False

def set_text_content(shape, text):
    """Set the text content of a shape."""
    if hasattr(shape, 'text_frame'):
        shape.text_frame.text = text
        return True
    return False

def replace_text(shape, find_text, replace_with):
    """Replace specific text in a shape with new text."""
    if hasattr(shape, 'text_frame'):
        original_text = shape.text_frame.text
        if find_text in original_text:
            new_text = original_text.replace(find_text, replace_with)
            shape.text_frame.text = new_text
            return True
    return False

def set_vertical_alignment(shape, alignment="middle"):
    """Set vertical text alignment in a shape (top, middle, bottom)."""
    if hasattr(shape, 'text_frame'):
        alignment_map = {
            "top": MSO_VERTICAL_ANCHOR.TOP,
            "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
            "bottom": MSO_VERTICAL_ANCHOR.BOTTOM
        }
        shape.text_frame.vertical_anchor = alignment_map.get(alignment.lower(), MSO_VERTICAL_ANCHOR.MIDDLE)
        return True
    return False

def set_autofit(shape, autofit_type="none"):
    """Set text autofit options (none, shape, text)."""
    if hasattr(shape, 'text_frame'):
        autofit_map = {
            "none": MSO_AUTO_SIZE.NONE,
            "shape": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
            "text": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        }
        shape.text_frame.auto_size = autofit_map.get(autofit_type.lower(), MSO_AUTO_SIZE.NONE)
        return True
    return False

# ----- BULLET POINT FUNCTIONS -----
def remove_bullets(shape):
    """Remove bullets from text in a shape."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            try:
                paragraph._p.get_or_add_pPr().remove_child('buAutoNum')
                paragraph._p.get_or_add_pPr().remove_child('buChar')
                paragraph._p.get_or_add_pPr().remove_child('buBlip')
            except:
                pass  
        return True
    return False

def add_bullets(shape, level=0):
    """Add bullets to text in a shape."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.level = level
        return True
    return False

def set_bullet_character(shape, bullet_char="•"):
    """Set custom bullet character for a shape."""
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.bullet_character = bullet_char
        return True
    return False


# ----- SHAPE MANIPULATION FUNCTIONS -----
def resize_shape(shape, width=Inches(2), height=Inches(1)):
    """Resize a shape (e.g., image, table)."""
    try:
        shape.width = width
        shape.height = height
        return True
    except Exception as e:
        logging.error(f"Error resizing shape: {e}")
        return False

def move_shape(shape, left=Inches(1), top=Inches(1)):
    """Move a shape to a new position."""
    try:
        shape.left = left
        shape.top = top
        return True
    except Exception as e:
        logging.error(f"Error moving shape: {e}")
        return False
    
def remove_shape(shape):
    """Remove a shape from a slide."""
    try:
        if shape._element is not None:
            shape._element.getparent().remove(shape._element)
            return True
        return False
    except Exception as e:
        logging.error(f"Error removing shape: {e}")
        return False

def set_shape_fill(shape, color_rgb=(255, 255, 255)):
    """Set the fill color of a shape."""
    try:
        fill = shape.fill
        fill.solid()
        if isinstance(color_rgb, list):
            color_rgb = tuple(color_rgb)
        fill.fore_color.rgb = RGBColor(*color_rgb)
        return True
    except Exception as e:
        logging.error(f"Error setting shape fill: {e}")
        return False

def set_shape_transparency(shape, transparency=0):
    """Set the transparency of a shape (0-100%)."""
    try:
        if hasattr(shape.fill, 'transparency'):
            shape.fill.transparency = min(max(transparency/100, 0), 1) 
            return True
        return False
    except Exception as e:
        logging.error(f"Error setting shape transparency: {e}")
        return False

def set_shape_border(shape, color_rgb=(0, 0, 0), width=Pt(1)):
    """Set the border color and width of a shape."""
    try:
        line = shape.line
        if isinstance(color_rgb, list):
            color_rgb = tuple(color_rgb)
        line.color.rgb = RGBColor(*color_rgb)
        line.width = width
        return True
    except Exception as e:
        logging.error(f"Error setting shape border: {e}")
        return False

def remove_shape_border(shape):
    """Remove the border from a shape."""
    try:
        line = shape.line
        line.fill.background()
        return True
    except Exception as e:
        logging.error(f"Error removing shape border: {e}")
        return False


# ----- SLIDE FUNCTIONS -----
def set_background_color(slide, color_rgb=(255, 255, 255)):
    """Set the slide background color."""
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        
        if isinstance(color_rgb, list):
            color_rgb = tuple(color_rgb)
        fill.fore_color.rgb = RGBColor(*color_rgb)
        return True
    except Exception as e:
        logging.error(f"Error setting background color: {e}")
        return False


# ----- TABLE FUNCTIONS -----
def standardize_table(slide, table_index=0, font_name="Arial", font_size=12):
    """Standardize a table's font and size."""
    try:
        table_shapes = []
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'table'):
                if i == table_index:
                    table_shapes.append(shape)
                    break
                
        if not table_shapes:
            return False
            
        table = table_shapes[0].table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.size = Pt(font_size)
        return True
    except Exception as e:
        logging.error(f"Error standardizing table: {e}")
        return False

def set_table_header(slide, table_index=0, header_color_rgb=(200, 200, 200)):
    """Format the header row of a table with background color."""
    try:
        table_shapes = []
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'table'):
                if len(table_shapes) == table_index:
                    table_shapes.append(shape)
                
        if not table_shapes:
            return False
            
        table = table_shapes[0].table
        if len(table.rows) > 0:
            header_row = table.rows[0]
            for cell in header_row.cells:
                if isinstance(header_color_rgb, list):
                    header_color_rgb = tuple(header_color_rgb)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*header_color_rgb)
            return True
        return False
    except Exception as e:
        logging.error(f"Error setting table header: {e}")
        return False

def set_table_cell_color(slide, table_index=0, row_index=0, col_index=0, color_rgb=(200, 200, 200)):
    """Set the background color of a specific table cell."""
    try:
        table_shapes = []
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'table'):
                if len(table_shapes) == table_index:
                    table_shapes.append(shape)
                
        if not table_shapes:
            return False
            
        table = table_shapes[0].table
        if len(table.rows) > row_index and len(table.columns) > col_index:
            cell = table.cell(row_index, col_index)
            if isinstance(color_rgb, list):
                color_rgb = tuple(color_rgb)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*color_rgb)
            return True
        return False
    except Exception as e:
        logging.error(f"Error setting table cell color: {e}")
        return False

def set_table_border(slide, table_index=0, border_width=Pt(1), border_color_rgb=(0, 0, 0)):
    """Set the border for all cells in a table."""
    try:
        table_shapes = []
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'table'):
                if len(table_shapes) == table_index:
                    table_shapes.append(shape)
                
        if not table_shapes:
            return False
            
        table = table_shapes[0].table
        for row in table.rows:
            for cell in row.cells:
                if not hasattr(cell, 'line'):
                    continue
                if isinstance(border_color_rgb, list):
                    border_color_rgb = tuple(border_color_rgb)
                cell.line.color.rgb = RGBColor(*border_color_rgb)
                cell.line.width = border_width
        return True
    except Exception as e:
        logging.error(f"Error setting table border: {e}")
        return False


# ----- FUNCTION MAP -----
function_map = {
    # Text formatting functions
    "set_font": set_font,
    "align_text": align_text,
    "adjust_spacing": adjust_spacing,
    "set_text_color": set_text_color,
    "apply_bold": apply_bold,
    "apply_italic": apply_italic,
    "remove_bold": remove_bold,
    "remove_italic": remove_italic,
    "set_text_content": set_text_content,
    "replace_text": replace_text,
    "set_vertical_alignment": set_vertical_alignment,
    "set_autofit": set_autofit,
    
    # Bullet point functions
    "add_bullets": add_bullets,
    "remove_bullets": remove_bullets,
    "set_bullet_character": set_bullet_character,
    
    # Shape manipulation functions
    "resize_shape": resize_shape,
    "move_shape": move_shape,
    "remove_shape": remove_shape,
    "set_shape_fill": set_shape_fill,
    "set_shape_transparency": set_shape_transparency,
    "set_shape_border": set_shape_border,
    "remove_shape_border": remove_shape_border,
    
    # Slide functions
    "set_background_color": set_background_color,
    
    # Table functions
    "standardize_table": standardize_table,
    "set_table_header": set_table_header,
    "set_table_cell_color": set_table_cell_color,
    "set_table_border": set_table_border,
}

# List of available predefined functions with signatures and descriptions
AVAILABLE_FUNCTIONS = [
    # Text formatting functions
    {
        "name": "set_font",
        "signature": "set_font(shape, font_name='Arial', font_size=12)",
        "description": "Changes the font name and size of text in a shape."
    },
    {
        "name": "align_text",
        "signature": "align_text(shape, alignment='center')",
        "description": "Aligns text in a shape (left, center, right)."
    },
    {
        "name": "adjust_spacing",
        "signature": "adjust_spacing(shape, line_spacing=1.0)",
        "description": "Adjusts line spacing in a shape's text frame."
    },
    {
        "name": "set_text_color",
        "signature": "set_text_color(shape, color_rgb=(0, 0, 0))",
        "description": "Sets text color in a shape using RGB values."
    },
    {
        "name": "apply_bold",
        "signature": "apply_bold(shape)",
        "description": "Applies bold formatting to text in a shape."
    },
    {
        "name": "apply_italic",
        "signature": "apply_italic(shape)",
        "description": "Applies italic formatting to text in a shape."
    },
    {
        "name": "remove_bold",
        "signature": "remove_bold(shape)",
        "description": "Removes bold formatting from text in a shape."
    },
    {
        "name": "remove_italic",
        "signature": "remove_italic(shape)",
        "description": "Removes italic formatting from text in a shape."
    },
    {
        "name": "set_text_content",
        "signature": "set_text_content(shape, text)",
        "description": "Sets the text content of a shape."
    },
    {
        "name": "replace_text",
        "signature": "replace_text(shape, find_text, replace_with)",
        "description": "Replaces specific text in a shape with new text."
    },
    {
        "name": "set_vertical_alignment",
        "signature": "set_vertical_alignment(shape, alignment='middle')",
        "description": "Sets vertical text alignment in a shape (top, middle, bottom)."
    },
    {
        "name": "set_autofit",
        "signature": "set_autofit(shape, autofit_type='none')",
        "description": "Sets text autofit options (none, shape, text)."
    },
    
    # Bullet point functions
    {
        "name": "remove_bullets",
        "signature": "remove_bullets(shape)",
        "description": "Removes bullets from text in a shape."
    },
    {
        "name": "add_bullets",
        "signature": "add_bullets(shape, level=0)",
        "description": "Adds bullets to text in a shape with specified level."
    },
    {
        "name": "set_bullet_character",
        "signature": "set_bullet_character(shape, bullet_char='•')",
        "description": "Sets custom bullet character for a shape."
    },
    
    # Shape manipulation functions
    {
        "name": "resize_shape",
        "signature": "resize_shape(shape, width=Inches(2), height=Inches(1))",
        "description": "Resizes a shape (e.g., image, table)."
    },
    {
        "name": "move_shape",
        "signature": "move_shape(shape, left=Inches(1), top=Inches(1))",
        "description": "Moves a shape to a new position."
    },
    {
        "name": "remove_shape",
        "signature": "remove_shape(shape)",
        "description": "Removes a shape from a slide."
    },
    {
        "name": "set_shape_fill",
        "signature": "set_shape_fill(shape, color_rgb=(255, 255, 255))",
        "description": "Sets the fill color of a shape."
    },
    {
        "name": "set_shape_transparency",
        "signature": "set_shape_transparency(shape, transparency=0)",
        "description": "Sets the transparency of a shape (0-100%)."
    },
    {
        "name": "set_shape_border",
        "signature": "set_shape_border(shape, color_rgb=(0, 0, 0), width=Pt(1))",
        "description": "Sets the border color and width of a shape."
    },
    {
        "name": "remove_shape_border",
        "signature": "remove_shape_border(shape)",
        "description": "Removes the border from a shape."
    },
    
    # Slide functions
    {
        "name": "set_background_color",
        "signature": "set_background_color(slide, color_rgb=(255, 255, 255))",
        "description": "Sets the slide background color."
    },

    
    # Table functions
    {
        "name": "standardize_table",
        "signature": "standardize_table(slide, table_index=0, font_name='Arial', font_size=12)",
        "description": "Standardizes a table's font and size."
    },
    {
        "name": "set_table_header",
        "signature": "set_table_header(slide, table_index=0, header_color_rgb=(200, 200, 200))",
        "description": "Formats the header row of a table with background color."
    },
    {
        "name": "set_table_cell_color",
        "signature": "set_table_cell_color(slide, table_index=0, row_index=0, col_index=0, color_rgb=(200, 200, 200))",
        "description": "Sets the background color of a specific table cell."
    },
    {
        "name": "set_table_border",
        "signature": "set_table_border(slide, table_index=0, border_width=Pt(1), border_color_rgb=(0, 0, 0))",
        "description": "Sets the border for all cells in a table."
    }
]


FUNCTION_MAPPING_PROMPT = """
You are an AI assistant tasked with mapping PowerPoint modification tasks to predefined Python functions. 
Given a task specification and a list of available functions, determine which function to call and extract the necessary arguments from the task's parameters.

## Available Functions:
{functions_list}

## Task Specification:
{task_json}

## Instructions:
1. Analyze the task's `action`, `task_description`, and `params` to determine the most appropriate function to call.
2. Extract or transform values from the task's `params` to match the function's signature.
3. If the task cannot be mapped to any function, return an empty JSON object (`{{}}`).
4. If the task can be mapped, output a JSON object with:
   - "function_name": the name of the function to call
   - "arguments": a dictionary of argument names and their values


## Examples:
### Example 1:
Task Specification:
{{
  "action": "change_font",
  "task_description": "Change the font to Arial and size to 14pt",
  "params": {{"font_name": "Arial", "font_size": 14}}
}}

Output:
{{
  "function_name": "set_font",
  "arguments": {{"font_name": "Arial", "font_size": 14}}
}}

### Example 2:
Task Specification:
{{
  "action": "change_text_color",
  "task_description": "Change the text color to orange",
  "params": {{"color": "orange"}}
}}

Output:
{{
  "function_name": "set_text_color",
  "arguments": {{"color_rgb": [255, 165, 0]}}
}}

### Example 3:
Task Specification:
{{
  "action": "unknown_action",
  "task_description": "Do something undefined",
  "params": {{}}
}}

Output:
{{}}

## CRITICAL INSTRUCTIONS:
- Only map if the function perfectly matches the task's intent AND you can extract ALL necessary arguments from the input task's params.
- Ensure argument types match the function signature and have the correct data type (string, int, float, boolean, list for RGB like [255, 255, 255] not colour name string) 
- If a parameter is missing, use the function's default value if applicable.
- For functions requiring a `shape` or `slide` argument, assume it will be provided separately and do not include it in the arguments.
- If the task is too complex or requires logic not covered by a single function, return an empty JSON object (`{{}}`).

Your response must be a valid JSON object.
"""

def get_llm_mapped_function(task):
    functions_list = json.dumps(AVAILABLE_FUNCTIONS, indent=2)
    task_json = json.dumps(task, indent=2)
    
    prompt = FUNCTION_MAPPING_PROMPT.format(functions_list=functions_list, task_json=task_json)
    
    try:
        response = client.models.generate_content(model="gemini-2.0-flash", contents=[prompt])
        logging.info(f"LLM function mapping response: {response.text}")
        json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
        if json_match:
            mapping = json.loads(json_match.group(0))
            logging.info(f"LLM function mapping result: {mapping}")
            if "function_name" in mapping and "arguments" in mapping:
                return mapping
        return {}
    except Exception as e:
        logging.error(f"Failed to get LLM mapping for task: {e}")
        return {}
    





def find_shape_by_hint(slide, hint):
    if not hint or not isinstance(hint, str) or not hint.strip():
        return slide.shapes[0] if len(slide.shapes) > 0 else None
    
    hint = hint.strip().lower()
    
    # Case 1: Shape by index (e.g., "shape:2")
    shape_index_match = re.match(r'shape:(\d+)', hint)
    if shape_index_match:
        index = int(shape_index_match.group(1))
        if 0 <= index < len(slide.shapes):
            return slide.shapes[index]
    
    # Case 2: Placeholder by index (e.g., "placeholder:0")
    placeholder_index_match = re.match(r'placeholder:(\d+)', hint)
    if placeholder_index_match:
        index = int(placeholder_index_match.group(1))
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder and shape.placeholder_format.idx == index:
                return shape
    
    # Case 3: Exact text match
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame.text.strip().lower() == hint:
            return shape
    
    # Case 4: Common placeholders by name
    placeholder_mapping = {
        'title': 0,
        'subtitle': 1,
        'content': 2,
        'body': 2,
        'text': 2,
        'footer': 15,
        'header': 16,
        'date': 17,
        'slide number': 18
    }
    
    for keyword, idx in placeholder_mapping.items():
        if keyword in hint:
            for shape in slide.shapes:
                if (hasattr(shape, 'is_placeholder') and 
                    shape.is_placeholder and 
                    shape.placeholder_format.idx == idx):
                    return shape
    
    # Case 5: Shape type search
    shape_type_hints = {
        'table': MSO_SHAPE_TYPE.TABLE,
        'chart': MSO_SHAPE_TYPE.CHART,
        'image': MSO_SHAPE_TYPE.PICTURE,
        'picture': MSO_SHAPE_TYPE.PICTURE,
        'text box': MSO_SHAPE_TYPE.TEXT_BOX,
        'text': MSO_SHAPE_TYPE.TEXT_BOX,
        'textbox': MSO_SHAPE_TYPE.TEXT_BOX,
        'group': MSO_SHAPE_TYPE.GROUP,
        'line': MSO_SHAPE_TYPE.LINE
    }
    
    for type_hint, shape_type in shape_type_hints.items():
        if type_hint in hint:
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type') and shape.shape_type == shape_type:
                    return shape
    
    # Case 6: Position-based search
    position_hints = {
        'top left': (0, 0, 0.5, 0.5),
        'top center': (0.25, 0, 0.75, 0.5),
        'top right': (0.5, 0, 1, 0.5),
        'middle left': (0, 0.25, 0.5, 0.75),
        'center': (0.25, 0.25, 0.75, 0.75),
        'middle right': (0.5, 0.25, 1, 0.75),
        'bottom left': (0, 0.5, 0.5, 1),
        'bottom center': (0.25, 0.5, 0.75, 1),
        'bottom right': (0.5, 0.5, 1, 1)
    }
    
    for pos_hint, (left_ratio, top_ratio, right_ratio, bottom_ratio) in position_hints.items():
        if pos_hint in hint:
            slide_width = slide.width
            slide_height = slide.height
            
            # Find shapes in this region
            for shape in slide.shapes:
                shape_left_ratio = shape.left / slide_width
                shape_top_ratio = shape.top / slide_height
                shape_right_ratio = (shape.left + shape.width) / slide_width
                shape_bottom_ratio = (shape.top + shape.height) / slide_height
                
                # Check if shape is mostly in the target region
                if (shape_left_ratio >= left_ratio and 
                    shape_top_ratio >= top_ratio and 
                    shape_right_ratio <= right_ratio and 
                    shape_bottom_ratio <= bottom_ratio):
                    return shape
    
    # Case 7: Partial text match
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and hint in shape.text_frame.text.lower():
            return shape
    
    # Fallback: Any text-containing shape
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
            return shape
    
    # Last resort: first shape
    return slide.shapes[0] if len(slide.shapes) > 0 else None
